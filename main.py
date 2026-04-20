import os
import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def remove_watermark_from_shapes(shapes, slide_width, slide_height, threshold_ratio=0.8):
    """
    Belirli bir şekil koleksiyonu içindeki sağ alt köşede bulunan
    resim nesnelerini tespit edip siler.
    """
    shapes_to_delete = []
    
    threshold_left = slide_width * threshold_ratio
    threshold_top = slide_height * threshold_ratio

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.left and shape.top:
                is_bottom_right = (shape.left > threshold_left) and (shape.top > threshold_top)
                
                if is_bottom_right:
                    shapes_to_delete.append(shape)

    # Tespit edilen şekilleri XML ağacından güvenli bir şekilde kaldır
    for shape in shapes_to_delete:
        element = shape._element
        element.getparent().remove(element)

def process_presentation(input_path, output_path):
    print(f"[*] İşlem başlatılıyor: {input_path}")
    
    if not os.path.exists(input_path):
        print("[-] Hata: Belirtilen girdi dosyası bulunamadı.")
        return

    try:
        prs = Presentation(input_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        print("[*] Slaytlar taranıyor...")
        for slide in prs.slides:
            remove_watermark_from_shapes(slide.shapes, slide_width, slide_height)
            
        print("[*] Slayt düzenleri (Layouts) taranıyor...")
        for layout in prs.slide_layouts:
            remove_watermark_from_shapes(layout.shapes, slide_width, slide_height)
            
        print("[*] Asıl Slayt (Slide Master) taranıyor...")
        for master_layout in prs.slide_master.slide_layouts:
            remove_watermark_from_shapes(master_layout.shapes, slide_width, slide_height)
        remove_watermark_from_shapes(prs.slide_master.shapes, slide_width, slide_height)

        prs.save(output_path)
        print(f"[+] İşlem başarılı! Temizlenen dosya: {output_path}")

    except Exception as e:
        print(f"[-] Slayt işlenirken bir hata oluştu: {str(e)}")

def main():
    parser = argparse.ArgumentParser(description="PPTX dosyalarından sağ alt köşedeki (örn. AI üretimi) filigran/logoları temizler.")
    parser.add_argument("-i", "--input", required=True, help="İşlenecek orijinal PPTX dosyasının yolu")
    parser.add_argument("-o", "--output", required=True, help="Temizlenmiş PPTX dosyasının kaydedileceği yol")
    
    args = parser.parse_args()
    process_presentation(args.input, args.output)

if __name__ == "__main__":
    main()
