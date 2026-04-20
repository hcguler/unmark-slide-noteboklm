#!/usr/bin/env python3
"""
remove_notebooklm_logo.py
━━━━━━━━━━━━━━━━━━━━━━━━━
Removes the NotebookLM watermark logo from all slides in a PowerPoint (.pptx) file.

The logo is embedded inside the background image of each slide (bottom-right corner).
This script detects it automatically via pixel analysis and fills it with the
surrounding background color — leaving the rest of the slide completely intact.

Usage:
    python remove_notebooklm_logo.py <input.pptx> [output.pptx]
    python remove_notebooklm_logo.py presentation.pptx
    python remove_notebooklm_logo.py presentation.pptx cleaned.pptx

Author: https://github.com/your-username/remove-notebooklm-logo
License: MIT
"""

import sys
import os
import io
import argparse
import logging
from pathlib import Path
from typing import Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
    import numpy as np
except ImportError:
    print("ERROR: Pillow or numpy not installed. Run: pip install Pillow numpy")
    sys.exit(1)


# ─── Configuration ────────────────────────────────────────────────────────────

DEFAULT_SEARCH_REGION = {
    "x_start_pct": 0.70,   # Search from 70% width to right edge
    "y_start_pct": 0.80,   # Search from 80% height to bottom edge
}

LOGO_DETECTION = {
    "bg_tolerance": 35,        # Max pixel deviation considered "background"
    "min_logo_pixels": 50,     # Min non-background pixels to confirm logo found
    "padding": 6,              # Extra pixels added around detected logo bounds
    "bg_sample_rows": 5,       # Rows above logo to sample background color
}

SUPPORTED_IMAGE_TYPES = {"image/png", "image/jpeg", "image/jpg"}

# ─── Logging ──────────────────────────────────────────────────────────────────

logging.basicConfig(
    level=logging.INFO,
    format="%(levelname)s  %(message)s"
)
log = logging.getLogger(__name__)


# ─── Core Logic ───────────────────────────────────────────────────────────────

def detect_background_color(arr: np.ndarray, logo_box: Tuple[int,int,int,int]) -> np.ndarray:
    """
    Sample the background color from the area just above the detected logo.
    Falls back to sampling from the bottom-right corner if needed.
    
    Args:
        arr: Image as numpy array (H, W, C)
        logo_box: (x1, y1, x2, y2) bounding box of detected logo
    
    Returns:
        Background color as numpy array [R, G, B, A] or [R, G, B]
    """
    x1, y1, x2, y2 = logo_box
    h, w = arr.shape[:2]
    rows = LOGO_DETECTION["bg_sample_rows"]

    # Priority 1: strip directly above logo
    if y1 > rows + 2:
        strip = arr[y1 - rows: y1, x1: x2 + 1]
        if strip.size > 0:
            return np.median(strip.reshape(-1, strip.shape[2]), axis=0).astype(np.uint8)

    # Priority 2: area directly below logo
    if y2 + rows < h:
        strip = arr[y2 + 1: y2 + rows + 1, x1: x2 + 1]
        if strip.size > 0:
            return np.median(strip.reshape(-1, strip.shape[2]), axis=0).astype(np.uint8)

    # Priority 3: far corner of the slide (should always be background)
    corner = arr[int(h * 0.97):h, int(w * 0.97):w]
    if corner.size > 0:
        return np.median(corner.reshape(-1, corner.shape[2]), axis=0).astype(np.uint8)

    # Fallback: last resort
    return arr[y1, x1].copy()


def estimate_background_from_corners(arr: np.ndarray) -> np.ndarray:
    """Estimate the dominant slide background color from corner samples."""
    h, w = arr.shape[:2]
    size = max(5, min(15, h // 50))
    corners = [
        arr[0:size, 0:size],
        arr[0:size, w - size:w],
        arr[h - size:h, 0:size],
        arr[h - size:h, w - size:w],
    ]
    all_pixels = np.concatenate([c.reshape(-1, arr.shape[2]) for c in corners], axis=0)
    return np.median(all_pixels, axis=0).astype(np.uint8)


def find_logo_bounds(arr: np.ndarray) -> Optional[Tuple[int,int,int,int]]:
    """
    Detect the NotebookLM logo in the bottom-right region of an image.
    
    Strategy: Any pixel that deviates significantly from the background color
    (within the search region) is considered part of the logo.
    
    Args:
        arr: Image as numpy array (H, W, channels)
    
    Returns:
        (x1, y1, x2, y2) bounding box with padding, or None if not found
    """
    h, w = arr.shape[:2]

    # Estimate background
    bg = estimate_background_from_corners(arr)

    # Define search region
    rx = int(w * DEFAULT_SEARCH_REGION["x_start_pct"])
    ry = int(h * DEFAULT_SEARCH_REGION["y_start_pct"])
    region = arr[ry:, rx:]

    # Find pixels that differ from background
    diff = np.abs(region[:, :, :3].astype(np.int32) - bg[:3].astype(np.int32))
    is_logo = diff.max(axis=2) > LOGO_DETECTION["bg_tolerance"]

    logo_rows = np.where(is_logo.any(axis=1))[0]
    logo_cols = np.where(is_logo.any(axis=0))[0]

    if len(logo_rows) < 1 or len(logo_cols) < 1:
        return None

    pixel_count = int(is_logo.sum())
    if pixel_count < LOGO_DETECTION["min_logo_pixels"]:
        log.debug(f"  Too few logo pixels ({pixel_count}), skipping")
        return None

    pad = LOGO_DETECTION["padding"]
    x1 = max(0, rx + int(logo_cols.min()) - pad)
    y1 = max(0, ry + int(logo_rows.min()) - pad)
    x2 = min(w - 1, rx + int(logo_cols.max()) + pad)
    y2 = min(h - 1, ry + int(logo_rows.max()) + pad)

    return (x1, y1, x2, y2)


def remove_logo_from_image(img_bytes: bytes, content_type: str) -> Tuple[bytes, bool]:
    """
    Remove the NotebookLM logo from image bytes.
    
    Args:
        img_bytes: Raw image bytes
        content_type: MIME type of image (e.g. 'image/png')
    
    Returns:
        (modified_image_bytes, was_logo_found)
    """
    fmt_map = {
        "image/png":  "PNG",
        "image/jpeg": "JPEG",
        "image/jpg":  "JPEG",
    }
    save_fmt = fmt_map.get(content_type, "PNG")

    img = Image.open(io.BytesIO(img_bytes))
    # Ensure RGBA for consistent processing
    original_mode = img.mode
    if img.mode not in ("RGB", "RGBA"):
        img = img.convert("RGBA")
    elif img.mode == "RGB":
        img = img.convert("RGBA")

    arr = np.array(img, dtype=np.uint8).copy()

    logo_box = find_logo_bounds(arr)
    if logo_box is None:
        return img_bytes, False

    x1, y1, x2, y2 = logo_box
    log.debug(f"  Logo box: ({x1},{y1}) → ({x2},{y2})")

    bg_color = detect_background_color(arr, logo_box)
    arr[y1: y2 + 1, x1: x2 + 1] = bg_color

    result_img = Image.fromarray(arr)

    # Convert back to original mode if needed
    if original_mode == "RGB" and save_fmt == "JPEG":
        result_img = result_img.convert("RGB")

    buf = io.BytesIO()
    save_kwargs = {"format": save_fmt}
    if save_fmt == "JPEG":
        save_kwargs["quality"] = 95
        save_kwargs["subsampling"] = 0
    result_img.save(buf, **save_kwargs)
    return buf.getvalue(), True


def process_presentation(input_path: str, output_path: str) -> dict:
    """
    Process a .pptx file and remove NotebookLM logos from all slides.
    
    Args:
        input_path: Path to the input .pptx file
        output_path: Path to save the cleaned .pptx file
    
    Returns:
        Summary dict with slides_processed, logos_removed, slides_unchanged
    """
    prs = Presentation(input_path)

    total_slides = len(prs.slides)
    logos_removed = 0
    slides_unchanged = 0

    # Track already-processed images (rId → new bytes) to avoid reprocessing
    # shared background images used on multiple slides
    processed_rels: dict = {}

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_logos = 0

        for shape in slide.shapes:
            if shape.shape_type != 13:  # 13 = PICTURE
                continue

            try:
                image = shape.image
            except AttributeError:
                continue

            content_type = image.content_type
            if content_type not in SUPPORTED_IMAGE_TYPES:
                log.debug(f"  Slide {slide_idx}: Skipping unsupported image type {content_type}")
                continue

            # Get relationship ID for this picture
            blip = shape._element.blipFill.blip
            ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            rId = blip.get(f"{{{ns}}}embed")

            if rId in processed_rels:
                # Already processed — reuse result
                new_bytes, found = processed_rels[rId]
                if found:
                    slide_logos += 1
                continue

            original_bytes = image.blob
            new_bytes, found = remove_logo_from_image(original_bytes, content_type)
            processed_rels[rId] = (new_bytes, found)

            if found:
                # Replace the image in the relationship
                image_part = slide.part.related_part(rId)
                image_part._blob = new_bytes
                slide_logos += 1
                log.info(f"  Slide {slide_idx}: Logo removed ✓")
            else:
                log.debug(f"  Slide {slide_idx}: No logo detected in image")

        if slide_logos == 0:
            slides_unchanged += 1

        logos_removed += slide_logos

    prs.save(output_path)

    return {
        "slides_processed": total_slides,
        "logos_removed": logos_removed,
        "slides_unchanged": slides_unchanged,
    }


# ─── CLI ──────────────────────────────────────────────────────────────────────

def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="remove_notebooklm_logo",
        description=(
            "Remove the NotebookLM watermark logo from all slides in a .pptx file.\n"
            "The logo is detected automatically via pixel analysis and erased cleanly."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python remove_notebooklm_logo.py presentation.pptx
  python remove_notebooklm_logo.py presentation.pptx cleaned.pptx
  python remove_notebooklm_logo.py slides/deck.pptx output/deck_clean.pptx
  python remove_notebooklm_logo.py deck.pptx --verbose
  python remove_notebooklm_logo.py deck.pptx --dry-run
        """,
    )
    parser.add_argument(
        "input",
        metavar="INPUT",
        help="Path to the input .pptx file",
    )
    parser.add_argument(
        "output",
        metavar="OUTPUT",
        nargs="?",
        default=None,
        help=(
            "Path for the cleaned output .pptx file. "
            "Defaults to <input_name>_cleaned.pptx in the same directory."
        ),
    )
    parser.add_argument(
        "-v", "--verbose",
        action="store_true",
        help="Show detailed per-slide debug output",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Detect logos and report without saving any file",
    )
    parser.add_argument(
        "--tolerance",
        type=int,
        default=LOGO_DETECTION["bg_tolerance"],
        metavar="N",
        help=f"Pixel color tolerance for logo detection (default: {LOGO_DETECTION['bg_tolerance']}). "
             "Increase for low-contrast logos.",
    )
    parser.add_argument(
        "--padding",
        type=int,
        default=LOGO_DETECTION["padding"],
        metavar="N",
        help=f"Extra pixels to erase around detected logo bounds (default: {LOGO_DETECTION['padding']})",
    )
    return parser


def resolve_output_path(input_path: str, output_arg: Optional[str]) -> str:
    if output_arg:
        return output_arg
    p = Path(input_path)
    return str(p.parent / f"{p.stem}_cleaned{p.suffix}")


def main():
    parser = build_parser()
    args = parser.parse_args()

    # Apply CLI overrides to config
    LOGO_DETECTION["bg_tolerance"] = args.tolerance
    LOGO_DETECTION["padding"] = args.padding

    if args.verbose:
        log.setLevel(logging.DEBUG)

    input_path = args.input
    if not os.path.isfile(input_path):
        log.error(f"Input file not found: {input_path}")
        sys.exit(1)
    if not input_path.lower().endswith(".pptx"):
        log.warning("Input file does not have .pptx extension — proceeding anyway.")

    output_path = resolve_output_path(input_path, args.output)

    if args.dry_run:
        log.info(f"DRY RUN — no file will be saved")
        log.info(f"Input:  {input_path}")
        # Run detection only (don't save)
        result = process_presentation(input_path, os.devnull)
        print(f"\n{'─'*50}")
        print(f"  DRY RUN RESULTS")
        print(f"  Slides scanned : {result['slides_processed']}")
        print(f"  Logos detected : {result['logos_removed']}")
        print(f"  No file written (dry-run mode)")
        print(f"{'─'*50}\n")
        return

    log.info(f"Input:  {input_path}")
    log.info(f"Output: {output_path}")
    log.info(f"Processing slides...\n")

    result = process_presentation(input_path, output_path)

    print(f"\n{'─'*50}")
    print(f"  DONE")
    print(f"  Slides processed : {result['slides_processed']}")
    print(f"  Logos removed    : {result['logos_removed']}")
    print(f"  Unchanged slides : {result['slides_unchanged']}")
    if result["logos_removed"] == 0:
        print(f"\n  ⚠  No logos were found. The file may already be clean,")
        print(f"     or try --tolerance 50 for a more aggressive search.")
    else:
        print(f"\n  ✓  Saved to: {output_path}")
    print(f"{'─'*50}\n")


if __name__ == "__main__":
    main()
